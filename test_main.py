import io
import json
from datetime import date
from unittest.mock import patch
from urllib.error import URLError

from fastapi.testclient import TestClient
from pptx import Presentation
from pptx.dml.color import RGBColor
from docx import Document
import openpyxl
from main import app
from utils.bible_api import BibleAPIError, generate_bible_text, normalize_korean_reference, parse_reference_lines
from utils.text_processing import (
    apply_combined_line_break,
    apply_equals_line_break,
    inspect_line_breaks,
    parse_bible_reference,
)

client = TestClient(app)


# ── GET / ────────────────────────────────────────────────────────────────────

def test_root():
    response = client.get('/')
    assert response.status_code == 200
    assert response.json() == {'message': 'Hello FastAPI'}


def test_client_page():
    response = client.get('/client')
    assert response.status_code == 200
    assert 'text/html' in response.headers['content-type']
    assert '<html lang="en">' in response.text
    assert 'Line Break Test Client' in response.text
    assert 'line break' in response.text
    assert 'bible lookup' in response.text
    assert '/line-break/combined/export_ppt' in response.text
    assert '/bible/generate' in response.text
    assert 'Generate Bible Text' in response.text
    assert 'Output copied.' in response.text
    assert 'Bible + Outline' in response.text
    assert 'General Text' in response.text
    assert 'id="combined-line-break-panel"' in response.text
    assert 'id="general-line-break-panel"' in response.text
    assert 'data-criterion="."' in response.text
    assert 'data-criterion=","' in response.text
    assert 'data-criterion="-"' in response.text
    assert 'Arabic number' in response.text
    assert 'Roman number' in response.text
    assert 'Chinese number' in response.text
    assert 'Bullet' in response.text
    assert 'Check mark' in response.text
    assert 'id="processGeneralBtn"' in response.text
    assert 'id="copyGeneralResultBtn"' in response.text
    assert 'id="copyCombinedResultBtn"' in response.text
    assert 'id="copyLookupBtn"' in response.text
    assert 'id="eraseCombinedSourceBtn"' in response.text
    assert 'id="eraseLookupInputBtn"' in response.text
    assert 'id="eraseResultBtn"' not in response.text
    assert 'id="eraseLookupResultBtn"' not in response.text
    assert 'height: 360px;' in response.text
    assert 'height: 260px;' in response.text
    assert 'resize: none;' in response.text
    assert 'overflow: auto;' in response.text
    assert 'scrollbar-width: none;' in response.text
    assert '-ms-overflow-style: none;' in response.text
    assert '::-webkit-scrollbar' in response.text
    assert 'display: none;' in response.text


def test_client_page_korean_query_locale():
    response = client.get('/client', params={'lang': 'ko'})
    assert response.status_code == 200
    assert '<html lang="ko">' in response.text
    assert '줄바꿈 테스트 클라이언트' in response.text
    assert '성경 본문 생성' in response.text
    assert '일반 텍스트' in response.text
    assert '기준' in response.text
    assert '아라비아 숫자' in response.text
    assert '체크 표시' in response.text
    assert '처리된 출력 복사' in response.text
    assert '일반 출력 복사' in response.text
    assert '출력이 복사되었습니다.' in response.text
    assert '/line-break/combined/export_ppt' in response.text
    assert '/bible/generate' in response.text


def test_client_page_korean_accept_language():
    response = client.get('/client', headers={'Accept-Language': 'ko-KR,ko;q=0.9,en;q=0.8'})
    assert response.status_code == 200
    assert '<html lang="ko">' in response.text
    assert '입력 텍스트' in response.text


def make_bible_payload(module: str, book_name: str, text: str):
    return {
        'results': [
            {
                'book_name': book_name,
                'verses': {
                    module: {
                        '3': {
                            '16': {
                                'chapter': 3,
                                'verse': 16,
                                'text': text,
                            }
                        }
                    }
                },
            }
        ]
    }


class MockHTTPResponse:
    def __init__(self, payload):
        self.payload = payload

    def __enter__(self):
        return io.StringIO(json.dumps(self.payload))

    def __exit__(self, exc_type, exc, tb):
        return False


class TestBibleLookup:
    @patch('utils.bible_api.urlopen')
    def test_korean_lookup(self, mock_urlopen):
        mock_urlopen.return_value = MockHTTPResponse(make_bible_payload('korean', '요한복음', '하나님이 세상을 이처럼 사랑하사'))
        response = client.get('/bible/lookup', params={'lang': 'ko', 'reference': '요 3:16'})
        assert response.status_code == 200
        data = response.json()
        assert data['lang'] == 'ko'
        assert data['module'] == 'korean'
        assert data['reference'] == '요 3:16'
        assert data['text'] == '하나님이 세상을 이처럼 사랑하사'
        assert data['verses'][0]['book'] == '요한복음'

    @patch('utils.bible_api.urlopen')
    def test_english_lookup(self, mock_urlopen):
        mock_urlopen.return_value = MockHTTPResponse(make_bible_payload('web', 'John', 'For God so loved the world'))
        response = client.get('/bible/lookup', params={'lang': 'en', 'reference': 'John 3:16'})
        assert response.status_code == 200
        data = response.json()
        assert data['module'] == 'web'
        assert data['text'] == 'For God so loved the world'

    @patch('utils.bible_api.urlopen')
    def test_japanese_lookup(self, mock_urlopen):
        mock_urlopen.return_value = MockHTTPResponse(make_bible_payload('kougo', 'ヨハネ', '神はそのひとり子をお与えになったほどに'))
        response = client.get('/bible/lookup', params={'lang': 'ja', 'reference': 'ヨハネ 3:16'})
        assert response.status_code == 200
        data = response.json()
        assert data['module'] == 'kougo'
        assert data['text'] == '神はそのひとり子をお与えになったほどに'

    def test_unsupported_language(self):
        response = client.get('/bible/lookup', params={'lang': 'fr', 'reference': 'Jean 3:16'})
        assert response.status_code == 400
        assert 'Unsupported language' in response.json()['detail']

    def test_missing_reference(self):
        response = client.get('/bible/lookup', params={'lang': 'en'})
        assert response.status_code == 400
        assert response.json()['detail'] == 'Reference is required.'

    def test_blank_reference(self):
        response = client.get('/bible/lookup', params={'lang': 'en', 'reference': '   '})
        assert response.status_code == 400
        assert response.json()['detail'] == 'Reference is required.'

    @patch('utils.bible_api.urlopen')
    def test_upstream_failure(self, mock_urlopen):
        mock_urlopen.side_effect = URLError('network down')
        response = client.get('/bible/lookup', params={'lang': 'en', 'reference': 'John 3:16'})
        assert response.status_code == 502
        assert response.json()['detail'] == 'Failed to fetch passage from Bible SuperSearch.'

    @patch('utils.bible_api.urlopen')
    def test_no_verses_found(self, mock_urlopen):
        mock_urlopen.return_value = MockHTTPResponse({'results': []})
        response = client.get('/bible/lookup', params={'lang': 'en', 'reference': 'Unknown 1:1'})
        assert response.status_code == 404
        assert response.json()['detail'] == 'No verses were returned for the given reference.'


class TestBibleReferenceParsing:
    def test_normalize_single_reference(self):
        parsed = normalize_korean_reference('마25:21')
        assert parsed['query_reference'] == '마태복음 25:21'
        assert parsed['display_reference'] == '마25:21'

    def test_normalize_range_reference(self):
        parsed = normalize_korean_reference('계2:8-13')
        assert parsed['query_reference'] == '요한계시록 2:8-13'
        assert parsed['display_reference'] == '계2:8-13'

    def test_normalize_pauline_reference(self):
        parsed = normalize_korean_reference('고전15:10')
        assert parsed['query_reference'] == '고린도전서 15:10'
        assert parsed['display_reference'] == '고전15:10'

    def test_parse_reference_lines_ignores_blank_and_deduplicates(self):
        input_count, parsed = parse_reference_lines('마25:21\n\n계2:8-13\n마25:21\n')
        assert input_count == 3
        assert [item['display_reference'] for item in parsed] == ['마25:21', '계2:8-13']

    def test_parse_reference_lines_extracts_single_parenthetical_reference(self):
        input_count, parsed = parse_reference_lines('(고전15:58)')
        assert input_count == 1
        assert [item['display_reference'] for item in parsed] == ['고전15:58']

    def test_parse_reference_lines_extracts_multiple_parenthetical_references(self):
        input_count, parsed = parse_reference_lines('(고전15:58)\n(대상6:31-32)')
        assert input_count == 2
        assert [item['display_reference'] for item in parsed] == ['고전15:58', '대상6:31-32']

    def test_parse_reference_lines_extracts_multiple_parens_on_one_line(self):
        input_count, parsed = parse_reference_lines('(고전15:58) (대상6:31-32)')
        assert input_count == 2
        assert [item['display_reference'] for item in parsed] == ['고전15:58', '대상6:31-32']


class TestBibleGeneration:
    @patch('utils.bible_api.fetch_bible_passage')
    def test_single_verse_generation(self, mock_fetch):
        mock_fetch.return_value = {
            'verses': [
                {'book': '마태복음', 'chapter': 25, 'verse': 21, 'text': '그 주인이 이르되'}
            ]
        }
        result = generate_bible_text('마25:21')
        assert result['input_count'] == 1
        assert result['unique_count'] == 1
        assert result['references'] == ['마25:21']
        assert result['output'] == '마25:21\n그 주인이 이르되'

    @patch('utils.bible_api.fetch_bible_passage')
    def test_range_generation_expands_per_verse(self, mock_fetch):
        mock_fetch.return_value = {
            'verses': [
                {'book': '요한계시록', 'chapter': 2, 'verse': 8, 'text': '서머나 교회의 사자에게'},
                {'book': '요한계시록', 'chapter': 2, 'verse': 9, 'text': '내가 네 환난과 궁핍을 아노니'},
            ]
        }
        result = generate_bible_text('계2:8-9')
        assert result['output'] == '계2:8\n서머나 교회의 사자에게\n\n계2:9\n내가 네 환난과 궁핍을 아노니'
        assert [item['output_reference'] for item in result['items']] == ['계2:8', '계2:9']

    @patch('utils.bible_api.fetch_bible_passage')
    def test_mixed_input_preserves_order_and_deduplicates(self, mock_fetch):
        def fake_fetch(lang, reference, timeout=15):
            mapping = {
                '마태복음 25:21': {
                    'verses': [{'book': '마태복음', 'chapter': 25, 'verse': 21, 'text': '잘 하였도다'}]
                },
                '이사야 43:7': {
                    'verses': [{'book': '이사야', 'chapter': 43, 'verse': 7, 'text': '내 이름으로 불려지는 모든 자'}]
                },
            }
            return mapping[reference]

        mock_fetch.side_effect = fake_fetch
        result = generate_bible_text('마25:21\n사43:7\n마25:21')
        assert result['input_count'] == 3
        assert result['unique_count'] == 2
        assert result['references'] == ['마25:21', '사43:7']
        assert result['output'] == '마25:21\n잘 하였도다\n\n사43:7\n내 이름으로 불려지는 모든 자'

    def test_invalid_reference_raises_400(self):
        response = client.post('/bible/generate', json={'text': '잘못된입력'})
        assert response.status_code == 400
        assert 'Unsupported or invalid reference' in response.json()['detail']

    @patch('utils.bible_api.fetch_bible_passage')
    def test_generate_endpoint_expands_range(self, mock_fetch):
        mock_fetch.return_value = {
            'verses': [
                {'book': '요한계시록', 'chapter': 2, 'verse': 8, 'text': '서머나 교회의 사자에게'},
                {'book': '요한계시록', 'chapter': 2, 'verse': 9, 'text': '내가 네 환난과 궁핍을 아노니'},
            ]
        }
        response = client.post('/bible/generate', json={'text': '계2:8-9'})
        assert response.status_code == 200
        data = response.json()
        assert data['references'] == ['계2:8-9']
        assert data['items'][0]['output_reference'] == '계2:8'
        assert data['items'][1]['output_reference'] == '계2:9'

    @patch('utils.bible_api.fetch_bible_passage')
    def test_generate_endpoint_upstream_failure(self, mock_fetch):
        mock_fetch.side_effect = BibleAPIError('Failed to fetch passage from Bible SuperSearch.')
        response = client.post('/bible/generate', json={'text': '마25:21'})
        assert response.status_code == 502


# ── utils.text_processing.inspect_line_breaks ────────────────────────────────

class TestInspectLineBreaks:
    def test_no_line_break(self):
        normalized, has_line_breaks, lines = inspect_line_breaks('창1:1 태초에 하나님이')
        assert normalized == '창1:1 태초에 하나님이'
        assert has_line_breaks is False
        assert lines == ['창1:1 태초에 하나님이']

    def test_unix_line_break(self):
        normalized, has_line_breaks, lines = inspect_line_breaks('창1:1\n태초에 하나님이')
        assert normalized == '창1:1\n태초에 하나님이'
        assert has_line_breaks is True
        assert lines == ['창1:1', '태초에 하나님이']

    def test_windows_line_break(self):
        normalized, has_line_breaks, lines = inspect_line_breaks('창1:1\r\n태초에 하나님이')
        assert normalized == '창1:1\n태초에 하나님이'
        assert has_line_breaks is True
        assert lines == ['창1:1', '태초에 하나님이']

    def test_old_mac_line_break(self):
        normalized, has_line_breaks, lines = inspect_line_breaks('창1:1\r태초에 하나님이')
        assert normalized == '창1:1\n태초에 하나님이'
        assert has_line_breaks is True
        assert lines == ['창1:1', '태초에 하나님이']


# ── POST /line-break ──────────────────────────────────────────────────────────

class TestLineBreak:
    def test_single_verse(self):
        response = client.post('/line-break', json={'text': '창1:1 태초에 하나님이 천지를 창조하시니라'})
        assert response.status_code == 200
        assert response.json()['result'] == '창1:1\n태초에 하나님이 천지를 창조하시니라'

    def test_multiple_verses(self):
        text = '마20:7 이르되 우리를 품꾼으로 쓰는 이가 없음이니이다\n\n마20:8 저물매 포도원 주인이'
        response = client.post('/line-break', json={'text': text})
        assert response.status_code == 200
        result = response.json()['result']
        assert '마20:7\n이르되' in result
        assert '\n\n마20:8\n' in result
        assert '마20:8\n저물매' in result

    def test_single_newline_normalized_to_double(self):
        text = '마20:7 이르되 우리를 품꾼으로 쓰는 이가 없음이니이다\n마20:8 저물매 포도원 주인이'
        response = client.post('/line-break', json={'text': text})
        assert response.status_code == 200
        result = response.json()['result']
        assert '\n\n마20:8\n' in result

    def test_triple_newline_preserved(self):
        text = '마25:21 착하고 충성된 종아\n\n\n마25:29 무릇 있는 자는'
        response = client.post('/line-break', json={'text': text})
        assert response.status_code == 200
        result = response.json()['result']
        assert '\n\n\n마25:29\n' in result

    def test_footnote_not_broken(self):
        text = '요1:1 태초에 1)말씀이 계시니라 이 1)말씀이 하나님과 함께 계셨으니'
        response = client.post('/line-break', json={'text': text})
        assert response.status_code == 200
        result = response.json()['result']
        assert '요1:1\n태초에' in result
        assert '1)말씀이' in result
        assert '1)\n말씀이' not in result

    def test_parenthetical_content_after_reference(self):
        text = '눅23:51 (그들의 결의와 행사에 찬성하지 아니한 자라) 그는 유대인의 동네 아리마대 사람이요 하나님의 나라를 기다리는 자라'
        response = client.post('/line-break', json={'text': text})
        assert response.status_code == 200
        assert response.json()['result'] == (
            '눅23:51\n'
            '(그들의 결의와 행사에 찬성하지 아니한 자라) 그는 유대인의 동네 아리마대 사람이요 하나님의 나라를 기다리는 자라'
        )

    def test_no_digit_korean_boundary(self):
        text = '태초에 하나님이 천지를 창조하시니라'
        response = client.post('/line-break', json={'text': text})
        assert response.status_code == 200
        assert response.json()['result'] == text

    def test_windows_newline_normalized(self):
        text = '마20:7 이르되 우리를 품꾼으로 쓰는 이가 없음이니이다\r\n마20:8 저물매 포도원 주인이'
        response = client.post('/line-break', json={'text': text})
        assert response.status_code == 200
        result = response.json()['result']
        assert '\r' not in result
        assert '\n\n마20:8\n' in result


# ── utils.text_processing.apply_equals_line_break ─────────────────────────────

class TestApplyEqualsLineBreak:
    def test_single_line_breaks_after_equals(self):
        text = '첫째:생활가난때=생활부요믿어야(고후8:9)'
        assert apply_equals_line_break(text) == (
            '첫째:생활가난때=\n생활부요믿어야(고후8:9)'
        )

    def test_multiple_lines_become_separate_blocks(self):
        text = (
            '첫째:생활가난때=생활부요믿어야(고후8:9)\n'
            '둘째:귀신역사때=귀신축귀믿어야(막16:16~17)'
        )
        assert apply_equals_line_break(text) == (
            '첫째:생활가난때=\n생활부요믿어야(고후8:9)\n\n'
            '둘째:귀신역사때=\n귀신축귀믿어야(막16:16-17)'
        )

    def test_multiple_equals_signs_each_break(self):
        text = '가=나=다'
        assert apply_equals_line_break(text) == '가=\n나=\n다'

    def test_blank_lines_are_dropped(self):
        text = '첫째:가난=부요(고후8:9)\n\n둘째:귀신=축귀(막16:17)'
        assert apply_equals_line_break(text) == (
            '첫째:가난=\n부요(고후8:9)\n\n둘째:귀신=\n축귀(막16:17)'
        )

    def test_line_without_equals_is_kept_as_is(self):
        assert apply_equals_line_break('그냥 문장입니다') == '그냥 문장입니다'

    def test_tilde_between_numbers_in_citation_becomes_dash(self):
        assert apply_equals_line_break('가나다=라마바(요1:1~2)') == '가나다=\n라마바(요1:1-2)'

    def test_long_content_breaks_before_citation(self):
        text = '첫째:아주길게이어지는내용입니다=이것도아주아주아주길게이어지는답변내용입니다(고후8:9)'
        result = apply_equals_line_break(text)
        assert '이것도아주아주아주길게이어지는답변내용입니다\n(고후8:9)' in result

    def test_short_content_does_not_break_before_citation(self):
        text = '첫째:생활가난때=생활부요믿어야(고후8:9)'
        assert '믿어야\n(고후8:9)' not in apply_equals_line_break(text)


# ── POST /line-break/equals ────────────────────────────────────────────────────

class TestLineBreakEquals:
    def test_single_line(self):
        response = client.post('/line-break/equals', json={'text': '첫째:생활가난때=생활부요믿어야(고후8:9)'})
        assert response.status_code == 200
        assert response.json()['result'] == '첫째:생활가난때=\n생활부요믿어야(고후8:9)'

    def test_two_lines_produce_two_blocks(self):
        text = (
            '첫째:생활가난때=생활부요믿어야(고후8:9)\n'
            '둘째:귀신역사때=귀신축귀믿어야(막16:16~17)'
        )
        response = client.post('/line-break/equals', json={'text': text})
        assert response.status_code == 200
        result = response.json()['result']
        assert '첫째:생활가난때=\n생활부요믿어야(고후8:9)' in result
        assert '둘째:귀신역사때=\n귀신축귀믿어야(막16:16-17)' in result
        assert result.count('\n\n') == 1


# ── POST /line-break/export_ppt ──────────────────────────────────────────────

class TestExportPPT:
    def test_status_and_content_type(self):
        response = client.post('/line-break/export_ppt', json={'text': '창1:1 태초에 하나님이'})
        assert response.status_code == 200
        assert 'presentationml' in response.headers['content-type']

    def test_content_disposition(self):
        response = client.post('/line-break/export_ppt', json={'text': '창1:1 태초에 하나님이'})
        assert f'{date.today():%Y-%m-%d}.pptx' in response.headers['content-disposition']

    def test_slide_text_content(self):
        response = client.post('/line-break/export_ppt', json={'text': '창1:1 태초에 하나님이 천지를 창조하시니라'})
        prs = Presentation(io.BytesIO(response.content))
        all_text = '\n'.join(
            shape.text_frame.text
            for slide in prs.slides
            for shape in slide.shapes
            if shape.has_text_frame
        )
        assert '창1:1' in all_text
        assert '태초에' in all_text

    def test_multiple_blocks_produce_multiple_slides(self):
        text = '창1:1 태초에 하나님이 천지를 창조하시니라\n\n창1:2 땅이 혼돈하고 공허하며'
        response = client.post('/line-break/export_ppt', json={'text': text})
        prs = Presentation(io.BytesIO(response.content))
        assert len(prs.slides) == 2

    def test_windows_double_newline_produces_multiple_slides(self):
        text = '창1:1 태초에 하나님이 천지를 창조하시니라\r\n\r\n창1:2 땅이 혼돈하고 공허하며'
        response = client.post('/line-break/export_ppt', json={'text': text})
        prs = Presentation(io.BytesIO(response.content))
        assert len(prs.slides) == 2

    def test_font_size_and_name_applied(self):
        response = client.post('/line-break/export_ppt', json={'text': '창1:1 태초에 하나님이'})
        prs = Presentation(io.BytesIO(response.content))
        run = prs.slides[0].shapes[0].text_frame.paragraphs[0].runs[0]
        assert run.font.size.pt == 52
        assert run.font.name == 'KoPubWorld바탕체 Bold'

    def test_east_asian_font_applied(self):
        from pptx.oxml.ns import qn
        response = client.post('/line-break/export_ppt', json={'text': '창1:1 태초에 하나님이'})
        prs = Presentation(io.BytesIO(response.content))
        run = prs.slides[0].shapes[0].text_frame.paragraphs[0].runs[0]
        ea = run.font._rPr.find(qn('a:ea'))
        assert ea is not None
        assert ea.get('typeface') == 'KoPubWorld바탕체 Bold'

    def test_slide_background_color_applied(self):
        response = client.post('/line-break/export_ppt', json={'text': '창1:1 태초에 하나님이'})
        prs = Presentation(io.BytesIO(response.content))
        background = prs.slides[0].background
        assert background.fill.fore_color.rgb == RGBColor(0x20, 0x38, 0x64)

    def test_font_color_applied(self):
        response = client.post('/line-break/export_ppt', json={'text': '창1:1 태초에 하나님이'})
        prs = Presentation(io.BytesIO(response.content))
        run = prs.slides[0].shapes[0].text_frame.paragraphs[0].runs[0]
        assert run.font.color.rgb == RGBColor(0xFF, 0xFF, 0xFF)

    def test_slide_size_is_widescreen(self):
        from pptx.util import Inches
        response = client.post('/line-break/export_ppt', json={'text': '창1:1 태초에 하나님이'})
        prs = Presentation(io.BytesIO(response.content))
        assert prs.slide_width == Inches(13.333)
        assert prs.slide_height == Inches(7.5)


# ── POST /line-break/equals/export_ppt ────────────────────────────────────────

class TestExportPPTEquals:
    def test_status_and_content_type(self):
        response = client.post('/line-break/equals/export_ppt', json={'text': '첫째:생활가난때=생활부요믿어야(고후8:9)'})
        assert response.status_code == 200
        assert 'presentationml' in response.headers['content-type']

    def test_content_disposition(self):
        response = client.post('/line-break/equals/export_ppt', json={'text': '첫째:생활가난때=생활부요믿어야(고후8:9)'})
        assert f'{date.today():%Y-%m-%d}.pptx' in response.headers['content-disposition']

    def test_slide_text_content(self):
        response = client.post('/line-break/equals/export_ppt', json={'text': '첫째:생활가난때=생활부요믿어야(고후8:9)'})
        prs = Presentation(io.BytesIO(response.content))
        all_text = '\n'.join(
            shape.text_frame.text
            for slide in prs.slides
            for shape in slide.shapes
            if shape.has_text_frame
        )
        assert '첫째:생활가난때=' in all_text
        assert '생활부요믿어야(고후8:9)' in all_text

    def test_each_input_line_produces_own_slide(self):
        text = (
            '첫째:생활가난때=생활부요믿어야(고후8:9)\n'
            '둘째:귀신역사때=귀신축귀믿어야(막16:16~17)'
        )
        response = client.post('/line-break/equals/export_ppt', json={'text': text})
        prs = Presentation(io.BytesIO(response.content))
        assert len(prs.slides) == 2

    def test_font_size_name_and_bold_applied(self):
        response = client.post('/line-break/equals/export_ppt', json={'text': '첫째:생활가난때=생활부요믿어야(고후8:9)'})
        prs = Presentation(io.BytesIO(response.content))
        run = prs.slides[0].shapes[0].text_frame.paragraphs[0].runs[0]
        assert run.font.size.pt == 60
        assert run.font.name == '맑은 고딕'
        assert run.font.bold is True

    def test_east_asian_font_applied(self):
        from pptx.oxml.ns import qn
        response = client.post('/line-break/equals/export_ppt', json={'text': '첫째:생활가난때=생활부요믿어야(고후8:9)'})
        prs = Presentation(io.BytesIO(response.content))
        run = prs.slides[0].shapes[0].text_frame.paragraphs[0].runs[0]
        ea = run.font._rPr.find(qn('a:ea'))
        assert ea is not None
        assert ea.get('typeface') == '맑은 고딕'

    def test_slide_background_color_is_black(self):
        response = client.post('/line-break/equals/export_ppt', json={'text': '첫째:생활가난때=생활부요믿어야(고후8:9)'})
        prs = Presentation(io.BytesIO(response.content))
        background = prs.slides[0].background
        assert background.fill.fore_color.rgb == RGBColor(0x00, 0x00, 0x00)

    def test_font_color_is_white(self):
        response = client.post('/line-break/equals/export_ppt', json={'text': '첫째:생활가난때=생활부요믿어야(고후8:9)'})
        prs = Presentation(io.BytesIO(response.content))
        run = prs.slides[0].shapes[0].text_frame.paragraphs[0].runs[0]
        assert run.font.color.rgb == RGBColor(0xFF, 0xFF, 0xFF)

    def test_slide_size_is_widescreen(self):
        from pptx.util import Inches
        response = client.post('/line-break/equals/export_ppt', json={'text': '첫째:생활가난때=생활부요믿어야(고후8:9)'})
        prs = Presentation(io.BytesIO(response.content))
        assert prs.slide_width == Inches(13.333)
        assert prs.slide_height == Inches(7.5)

    def test_text_alignment_is_centered(self):
        from pptx.enum.text import PP_ALIGN
        text = (
            '첫째:생활가난때=생활부요믿어야(고후8:9)\n'
            '둘째:귀신역사때=귀신축귀믿어야(막16:16~17)'
        )
        response = client.post('/line-break/equals/export_ppt', json={'text': text})
        prs = Presentation(io.BytesIO(response.content))
        for slide in prs.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    assert paragraph.alignment == PP_ALIGN.CENTER

    def test_text_box_is_centered_in_slide(self):
        from pptx.enum.text import MSO_ANCHOR
        from pptx.util import Inches
        response = client.post('/line-break/equals/export_ppt', json={'text': '첫째:생활가난때=생활부요믿어야(고후8:9)'})
        prs = Presentation(io.BytesIO(response.content))
        shape = prs.slides[0].shapes[0]
        assert shape.text_frame.vertical_anchor == MSO_ANCHOR.MIDDLE
        left_margin = shape.left
        right_margin = prs.slide_width - (shape.left + shape.width)
        top_margin = shape.top
        bottom_margin = prs.slide_height - (shape.top + shape.height)
        assert left_margin == right_margin
        assert top_margin == bottom_margin

    def test_align_defaults_to_center(self):
        from pptx.enum.text import PP_ALIGN
        response = client.post('/line-break/equals/export_ppt', json={'text': '첫째:생활가난때=생활부요믿어야(고후8:9)'})
        prs = Presentation(io.BytesIO(response.content))
        run = prs.slides[0].shapes[0].text_frame.paragraphs[0]
        assert run.alignment == PP_ALIGN.CENTER

    def test_align_left(self):
        from pptx.enum.text import PP_ALIGN
        response = client.post('/line-break/equals/export_ppt', json={'text': '첫째:생활가난때=생활부요믿어야(고후8:9)', 'align': 'left'})
        prs = Presentation(io.BytesIO(response.content))
        paragraph = prs.slides[0].shapes[0].text_frame.paragraphs[0]
        assert paragraph.alignment == PP_ALIGN.LEFT

    def test_align_right(self):
        from pptx.enum.text import PP_ALIGN
        response = client.post('/line-break/equals/export_ppt', json={'text': '첫째:생활가난때=생활부요믿어야(고후8:9)', 'align': 'right'})
        prs = Presentation(io.BytesIO(response.content))
        paragraph = prs.slides[0].shapes[0].text_frame.paragraphs[0]
        assert paragraph.alignment == PP_ALIGN.RIGHT

    def test_align_invalid_value_rejected(self):
        response = client.post('/line-break/equals/export_ppt', json={'text': '첫째:생활가난때=생활부요믿어야(고후8:9)', 'align': 'justify'})
        assert response.status_code == 422


# ── utils.text_processing.apply_combined_line_break ───────────────────────────

class TestApplyCombinedLineBreak:
    def test_equals_line_is_classified_equals(self):
        entries = apply_combined_line_break('첫째:날마다=그리스도말씀들어야(롬10:17)')
        assert entries == [('equals', ['첫째:날마다=\n그리스도말씀들어야(롬10:17)'])]

    def test_bible_line_is_classified_bible(self):
        text = '롬10:17 그러므로 믿음은 들음에서 나며 들음은 그리스도의 말씀으로 말미암았느니라'
        entries = apply_combined_line_break(text)
        assert entries == [
            ('bible', ['롬10:17\n그러므로 믿음은 들음에서 나며 들음은 그리스도의 말씀으로 말미암았느니라'])
        ]

    def test_mixed_lines_preserve_order_and_kind(self):
        text = (
            '첫째:날마다=그리스도말씀들어야(롬10:17)\n'
            '롬10:17 그러므로 믿음은 들음에서 나며 들음은 그리스도의 말씀으로 말미암았느니라'
        )
        entries = apply_combined_line_break(text)
        assert entries == [
            ('equals', ['첫째:날마다=\n그리스도말씀들어야(롬10:17)']),
            ('bible', ['롬10:17\n그러므로 믿음은 들음에서 나며 들음은 그리스도의 말씀으로 말미암았느니라']),
        ]

    def test_blank_lines_are_dropped(self):
        text = '첫째:날마다=그리스도말씀들어야(롬10:17)\n\n롬10:17 그러므로 믿음은 들음에서 나며'
        entries = apply_combined_line_break(text)
        assert len(entries) == 2

    def test_verses_listed_before_outline_are_reordered_next_to_citation(self):
        # Bible verses pasted up front, outline written after: the outline
        # entry must be paired with its cited verse, not left in place.
        text = (
            '고후8:9 우리 주 예수 그리스도의 은혜를 너희가 알거니와\n\n'
            '첫째:생활가난때=생활부요믿어야(고후8:9)'
        )
        entries = apply_combined_line_break(text)
        assert entries == [
            ('equals', ['첫째:생활가난때=\n생활부요믿어야(고후8:9)']),
            ('bible', ['고후8:9\n우리 주 예수 그리스도의 은혜를 너희가 알거니와']),
        ]

    def test_multi_verse_group_single_blank_line_stays_together_as_separate_parts(self):
        text = (
            '막16:16 믿고 세례를 받는 사람은 구원을 얻을 것이요\n\n'
            '막16:17 믿는 자들에게는 이런 표적이 따르리니\n\n'
            '둘째:귀신역사때=귀신축귀믿어야(막16:16~17)'
        )
        entries = apply_combined_line_break(text)
        assert entries == [
            ('equals', ['둘째:귀신역사때=\n귀신축귀믿어야(막16:16-17)']),
            ('bible', [
                '막16:16\n믿고 세례를 받는 사람은 구원을 얻을 것이요',
                '막16:17\n믿는 자들에게는 이런 표적이 따르리니',
            ]),
        ]

    def test_double_blank_line_starts_a_new_verse_group(self):
        text = (
            '막16:16 믿고 세례를 받는 사람은 구원을 얻을 것이요\n\n\n'
            '막16:17 믿는 자들에게는 이런 표적이 따르리니'
        )
        entries = apply_combined_line_break(text)
        assert entries == [
            ('bible', ['막16:16\n믿고 세례를 받는 사람은 구원을 얻을 것이요']),
            ('bible', ['막16:17\n믿는 자들에게는 이런 표적이 따르리니']),
        ]

    def test_header_lines_are_dropped(self):
        text = (
            '♡본론\n'
            '첫째:생활가난때=생활부요믿어야(고후8:9)\n'
            '♡결론\n'
            '둘째:날마다=성령충만해야(엡5:16)'
        )
        entries = apply_combined_line_break(text)
        assert entries == [
            ('equals', ['첫째:생활가난때=\n생활부요믿어야(고후8:9)']),
            ('equals', ['둘째:날마다=\n성령충만해야(엡5:16)']),
        ]

    def test_unmatched_citation_yields_equals_only(self):
        entries = apply_combined_line_break('첫째:생활가난때=생활부요믿어야(고후8:9)')
        assert entries == [('equals', ['첫째:생활가난때=\n생활부요믿어야(고후8:9)'])]

    def test_numbered_dot_marker_is_plain_text(self):
        entries = apply_combined_line_break('1.생활가난때=생활부요믿어야(고후8:9)')
        assert entries == [('equals', ['1.생활가난때=\n생활부요믿어야(고후8:9)'])]

    def test_numbered_paren_marker_is_plain_text(self):
        entries = apply_combined_line_break('1)생활가난때=생활부요믿어야(고후8:9)')
        assert entries == [('equals', ['1)생활가난때=\n생활부요믿어야(고후8:9)'])]

    def test_line_with_equals_but_no_ordinal_marker_is_dropped(self):
        # A bare '=' is no longer sufficient on its own; only lines that
        # start with an ordinal/list marker count as plain text.
        entries = apply_combined_line_break('참고: A=B 공식이 적용됨')
        assert entries == []

    def test_header_with_stray_equals_is_still_dropped(self):
        text = '♡본론 (수정=완료)\n첫째:생활가난때=생활부요믿어야(고후8:9)'
        entries = apply_combined_line_break(text)
        assert entries == [('equals', ['첫째:생활가난때=\n생활부요믿어야(고후8:9)'])]


# ── POST /line-break/combined ──────────────────────────────────────────────────

SERMON_TEXT = (
    '고후8:9 우리 주 예수 그리스도의 은혜를 너희가 알거니와 부요하신 이로서 너희를 위하여 가난하게 되심은 그의 가난함으로 말미암아 너희를 부요하게 하려 하심이라\n'
    '\n\n'
    '막16:16 믿고 세례를 받는 사람은 구원을 얻을 것이요 믿지 않는 사람은 정죄를 받으리라\n'
    '\n'
    '막16:17 믿는 자들에게는 이런 표적이 따르리니 곧 그들이 내 이름으로 귀신을 쫓아내며 새 방언을 말하며\n'
    '\n\n'
    '요16:33 이것을 너희에게 이르는 것은 너희로 내 안에서 평안을 누리게 하려 함이라 세상에서는 너희가 환난을 당하나 담대하라 내가 세상을 이기었노라\n'
    '\n\n'
    '막11:24 그러므로 내가 너희에게 말하노니 무엇이든지 기도하고 구하는 것은 받은 줄로 믿으라 그리하면 너희에게 그대로 되리라\n'
    '\n\n'
    '롬10:17 그러므로 믿음은 들음에서 나며 들음은 그리스도의 말씀으로 말미암았느니라\n'
    '\n\n'
    '엡5:16 세월을 아끼라 때가 악하니라\n'
    '\n'
    '엡5:17 그러므로 어리석은 자가 되지 말고 오직 주의 뜻이 무엇인가 이해하라\n'
    '\n'
    '엡5:18 술 취하지 말라 이는 방탕한 것이니 오직 성령으로 충만함을 받으라\n'
    '\n\n'
    '마4:10 이에 예수께서 말씀하시되 사탄아 물러가라 기록되었으되 주 너의 하나님께 경배하고 다만 그를 섬기라 하였느니라\n'
    '\n'
    '마4:11 이에 마귀는 예수를 떠나고 천사들이 나아와서 수종드니라\n'
    '\n'
    '♡본론\n'
    '첫째:생활가난때=생활부요믿어야(고후8:9)\n'
    '둘째:귀신역사때=귀신축귀믿어야(막16:16~17)\n'
    '셋째:환난인생때=환난이김믿어야(요16:33)\n'
    '넷째:기도간청때=기도응답믿어야(막11:24)\n'
    '♡결론\n'
    '첫째:날마다=그리스도말씀들어야(롬10:17)\n'
    '둘째:날마다=성령님으로충만해야(엡5:16~18)\n'
    '셋째:날마다=믿음의적사탄을대적(마4:10~11)'
)


class TestLineBreakCombined:
    def test_reference_only_line_consumes_following_verse_text(self):
        text = '마16:16\n주는 그리스도시요 살아 계신 하나님의 아들이시니이다'
        assert apply_combined_line_break(text) == [
            ('bible', ['마16:16\n주는 그리스도시요 살아 계신 하나님의 아들이시니이다'])
        ]

    def test_reference_only_line_is_preserved_at_end_of_input(self):
        assert apply_combined_line_break('마1:1') == [('bible', ['마1:1'])]

    def test_book_and_chapter_must_exist_in_catalog(self):
        assert apply_combined_line_break('회의1:2 다음 주 일정을 정합니다') == []
        assert apply_combined_line_break('막99:1 존재하지 않는 장입니다') == []

    def test_full_book_name_and_spacing_are_normalized(self):
        text = '마태복음 16장 16절 주는 그리스도시요'
        assert apply_combined_line_break(text) == [
            ('bible', ['마16:16\n주는 그리스도시요'])
        ]

    def test_catalog_key_variants_use_canonical_abbreviations(self):
        assert parse_bible_reference('눅2:1')['canonical_reference'] == '눅2:1'
        assert parse_bible_reference('약5:1')['canonical_reference'] == '약5:1'
        assert parse_bible_reference('합3:1')['canonical_reference'] == '합3:1'

    def test_range_separator_is_normalized_for_citation_pairing(self):
        text = (
            '첫째:믿음=고백해야(마태복음 16:16–17)\n'
            '마16:16 주는 그리스도시요\n'
            '마16:17 예수께서 대답하여 이르시되'
        )
        result = apply_combined_line_break(text)
        assert result[0][0] == 'equals'
        assert result[1] == (
            'bible',
            ['마16:16\n주는 그리스도시요', '마16:17\n예수께서 대답하여 이르시되'],
        )

    def test_mixed_input(self):
        text = (
            '첫째:날마다=그리스도말씀들어야(롬10:17)\n'
            '롬10:17 그러므로 믿음은 들음에서 나며 들음은 그리스도의 말씀으로 말미암았느니라'
        )
        response = client.post('/line-break/combined', json={'text': text})
        assert response.status_code == 200
        result = response.json()['result']
        assert '첫째:날마다=\n그리스도말씀들어야(롬10:17)' in result
        assert '롬10:17\n그러므로 믿음은 들음에서 나며 들음은 그리스도의 말씀으로 말미암았느니라' in result
        assert result.count('\n\n') == 1

    def test_sermon_manuscript_reorders_verses_next_to_their_outline_point(self):
        response = client.post('/line-break/combined', json={'text': SERMON_TEXT})
        assert response.status_code == 200
        result = response.json()['result']

        assert '♡본론' not in result
        assert '♡결론' not in result

        assert (
            '첫째:생활가난때=\n생활부요믿어야(고후8:9)\n\n\n'
            '고후8:9\n우리 주 예수 그리스도의 은혜를 너희가 알거니와 부요하신 이로서 너희를 위하여 가난하게 되심은 그의 가난함으로 말미암아 너희를 부요하게 하려 하심이라'
        ) in result
        assert (
            '둘째:귀신역사때=\n귀신축귀믿어야(막16:16-17)\n\n\n'
            '막16:16\n믿고 세례를 받는 사람은 구원을 얻을 것이요 믿지 않는 사람은 정죄를 받으리라\n\n'
            '막16:17\n믿는 자들에게는 이런 표적이 따르리니 곧 그들이 내 이름으로 귀신을 쫓아내며 새 방언을 말하며'
        ) in result
        assert (
            '둘째:날마다=\n성령님으로충만해야(엡5:16-18)\n\n\n'
            '엡5:16\n세월을 아끼라 때가 악하니라\n\n'
            '엡5:17\n그러므로 어리석은 자가 되지 말고 오직 주의 뜻이 무엇인가 이해하라\n\n'
            '엡5:18\n술 취하지 말라 이는 방탕한 것이니 오직 성령으로 충만함을 받으라'
        ) in result

        # Order: each outline point immediately precedes its own cited verse(s).
        assert result.index('첫째:생활가난때=') < result.index('고후8:9')
        assert result.index('고후8:9') < result.index('둘째:귀신역사때=')
        assert result.index('셋째:날마다=') < result.index('마4:10')


# ── POST /line-break/combined/export_ppt ──────────────────────────────────────

class TestExportPPTCombined:
    MIXED_TEXT = (
        '첫째:날마다=그리스도말씀들어야(롬10:17)\n'
        '롬10:17 그러므로 믿음은 들음에서 나며 들음은 그리스도의 말씀으로 말미암았느니라'
    )

    def test_status_and_content_type(self):
        response = client.post('/line-break/combined/export_ppt', json={'text': self.MIXED_TEXT})
        assert response.status_code == 200
        assert 'presentationml' in response.headers['content-type']

    def test_produces_one_slide_per_line(self):
        response = client.post('/line-break/combined/export_ppt', json={'text': self.MIXED_TEXT})
        prs = Presentation(io.BytesIO(response.content))
        assert len(prs.slides) == 2

    def test_equals_slide_keeps_outline_styling(self):
        from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
        response = client.post('/line-break/combined/export_ppt', json={'text': self.MIXED_TEXT})
        prs = Presentation(io.BytesIO(response.content))
        slide = prs.slides[0]
        shape = slide.shapes[0]
        run = shape.text_frame.paragraphs[0].runs[0]
        assert slide.background.fill.fore_color.rgb == RGBColor(0x00, 0x00, 0x00)
        assert run.font.name == '맑은 고딕'
        assert run.font.size.pt == 60
        assert run.font.bold is True
        assert run.font.color.rgb == RGBColor(0xFF, 0xFF, 0xFF)
        assert shape.text_frame.paragraphs[0].alignment == PP_ALIGN.CENTER
        assert shape.text_frame.vertical_anchor == MSO_ANCHOR.MIDDLE

    def test_bible_slide_keeps_bible_styling(self):
        response = client.post('/line-break/combined/export_ppt', json={'text': self.MIXED_TEXT})
        prs = Presentation(io.BytesIO(response.content))
        slide = prs.slides[1]
        shape = slide.shapes[0]
        run = shape.text_frame.paragraphs[0].runs[0]
        assert slide.background.fill.fore_color.rgb == RGBColor(0x20, 0x38, 0x64)
        assert run.font.name == 'KoPubWorld바탕체 Bold'
        assert run.font.size.pt == 52
        assert run.font.bold is None
        assert run.font.color.rgb == RGBColor(0xFF, 0xFF, 0xFF)
        assert shape.text_frame.paragraphs[0].alignment is None

    def test_align_applies_only_to_equals_slide(self):
        from pptx.enum.text import PP_ALIGN
        response = client.post('/line-break/combined/export_ppt', json={'text': self.MIXED_TEXT, 'align': 'left'})
        prs = Presentation(io.BytesIO(response.content))
        assert prs.slides[0].shapes[0].text_frame.paragraphs[0].alignment == PP_ALIGN.LEFT
        assert prs.slides[1].shapes[0].text_frame.paragraphs[0].alignment is None

    def test_sermon_manuscript_gives_every_verse_its_own_slide(self):
        # Multi-verse citations (막16:16~17, 엡5:16~18, 마4:10~11) must each
        # expand into one slide per verse rather than sharing a single slide.
        response = client.post('/line-break/combined/export_ppt', json={'text': SERMON_TEXT})
        prs = Presentation(io.BytesIO(response.content))
        assert len(prs.slides) == 18

        expected_kind_and_text = [
            ('equals', '첫째:생활가난때='),
            ('bible', '고후8:9'),
            ('equals', '둘째:귀신역사때='),
            ('bible', '막16:16'),
            ('bible', '막16:17'),
            ('equals', '셋째:환난인생때='),
            ('bible', '요16:33'),
            ('equals', '넷째:기도간청때='),
            ('bible', '막11:24'),
            ('equals', '첫째:날마다='),
            ('bible', '롬10:17'),
            ('equals', '둘째:날마다='),
            ('bible', '엡5:16'),
            ('bible', '엡5:17'),
            ('bible', '엡5:18'),
            ('equals', '셋째:날마다='),
            ('bible', '마4:10'),
            ('bible', '마4:11'),
        ]
        bible_bg = RGBColor(0x20, 0x38, 0x64)
        equals_bg = RGBColor(0x00, 0x00, 0x00)

        for slide, (kind, expected_prefix) in zip(prs.slides, expected_kind_and_text):
            shape = slide.shapes[0]
            assert shape.text_frame.text.startswith(expected_prefix)
            expected_bg = equals_bg if kind == 'equals' else bible_bg
            assert slide.background.fill.fore_color.rgb == expected_bg
            # A bible slide holds exactly one verse now, not a bundled group.
            if kind == 'bible':
                assert shape.text_frame.text.count('\n\n') == 0


# ── POST /line-break/export_docx ─────────────────────────────────────────────

class TestExportDOCX:
    def test_status_and_content_type(self):
        response = client.post('/line-break/export_docx', json={'text': '창1:1 태초에 하나님이'})
        assert response.status_code == 200
        assert 'wordprocessingml' in response.headers['content-type']

    def test_content_disposition(self):
        response = client.post('/line-break/export_docx', json={'text': '창1:1 태초에 하나님이'})
        assert 'output.docx' in response.headers['content-disposition']

    def test_paragraph_content(self):
        response = client.post('/line-break/export_docx', json={'text': '창1:1 태초에 하나님이 천지를 창조하시니라'})
        doc = Document(io.BytesIO(response.content))
        paragraphs = [p.text for p in doc.paragraphs]
        assert '창1:1' in paragraphs
        assert '태초에 하나님이 천지를 창조하시니라' in paragraphs

    def test_each_line_is_own_paragraph(self):
        text = '창1:1 태초에 하나님이\n창1:2 땅이 혼돈하고'
        response = client.post('/line-break/export_docx', json={'text': text})
        doc = Document(io.BytesIO(response.content))
        paragraphs = [p.text for p in doc.paragraphs]
        assert '창1:1' in paragraphs
        assert '태초에 하나님이' in paragraphs
        assert '창1:2' in paragraphs
        assert '땅이 혼돈하고' in paragraphs

    def test_windows_line_breaks_become_paragraphs(self):
        text = '창1:1 태초에 하나님이\r\n창1:2 땅이 혼돈하고'
        response = client.post('/line-break/export_docx', json={'text': text})
        doc = Document(io.BytesIO(response.content))
        paragraphs = [p.text for p in doc.paragraphs]
        assert '창1:1' in paragraphs
        assert '태초에 하나님이' in paragraphs
        assert '창1:2' in paragraphs
        assert '땅이 혼돈하고' in paragraphs


# ── POST /line-break/export_xlsx ─────────────────────────────────────────────

class TestExportXLSX:
    def test_status_and_content_type(self):
        response = client.post('/line-break/export_xlsx', json={'text': '창1:1 태초에 하나님이'})
        assert response.status_code == 200
        assert 'spreadsheetml' in response.headers['content-type']

    def test_content_disposition(self):
        response = client.post('/line-break/export_xlsx', json={'text': '창1:1 태초에 하나님이'})
        assert 'output.xlsx' in response.headers['content-disposition']

    def test_cell_content(self):
        response = client.post('/line-break/export_xlsx', json={'text': '창1:1 태초에 하나님이 천지를 창조하시니라'})
        wb = openpyxl.load_workbook(io.BytesIO(response.content))
        ws = wb.active
        cell_values = [row[0].value for row in ws.iter_rows()]
        assert '창1:1' in cell_values
        assert '태초에 하나님이 천지를 창조하시니라' in cell_values

    def test_each_line_in_own_row(self):
        text = '창1:1 태초에 하나님이\n창1:2 땅이 혼돈하고'
        response = client.post('/line-break/export_xlsx', json={'text': text})
        wb = openpyxl.load_workbook(io.BytesIO(response.content))
        ws = wb.active
        cell_values = [row[0].value for row in ws.iter_rows()]
        assert '창1:1' in cell_values
        assert '태초에 하나님이' in cell_values
        assert '창1:2' in cell_values
        assert '땅이 혼돈하고' in cell_values

    def test_windows_line_breaks_become_rows(self):
        text = '창1:1 태초에 하나님이\r\n창1:2 땅이 혼돈하고'
        response = client.post('/line-break/export_xlsx', json={'text': text})
        wb = openpyxl.load_workbook(io.BytesIO(response.content))
        ws = wb.active
        cell_values = [row[0].value for row in ws.iter_rows()]
        assert '창1:1' in cell_values
        assert '태초에 하나님이' in cell_values
        assert '창1:2' in cell_values
        assert '땅이 혼돈하고' in cell_values
