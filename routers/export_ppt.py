import io
from fastapi import APIRouter
from fastapi.responses import StreamingResponse
from pptx import Presentation
from pptx.util import Inches
from schemas.line_break import LineBreakRequest
from utils.pptx_helpers import pptx_download_filename, set_east_asian_font
from utils.slide_styles import BIBLE_STYLE, SLIDE_HEIGHT, SLIDE_WIDTH
from utils.text_processing import apply_line_break, inspect_line_breaks

router = APIRouter(prefix="/line-break/export_ppt", tags=["export"])

FONT_NAME = BIBLE_STYLE["font_name"]
FONT_SIZE = BIBLE_STYLE["font_size"]
FONT_COLOR = BIBLE_STYLE["font_color"]
BACKGROUND_COLOR = BIBLE_STYLE["background_color"]


@router.post('')
def export_ppt(request: LineBreakRequest):
    text = apply_line_break(request.text)
    text, _, _ = inspect_line_breaks(text)
    blocks = [b.strip() for b in text.split('\n\n') if b.strip()]
    if not blocks:
        blocks = [text.strip()]

    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    blank_layout = prs.slide_layouts[6]

    for block in blocks:
        slide = prs.slides.add_slide(blank_layout)
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = BACKGROUND_COLOR

        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), SLIDE_WIDTH - Inches(1), Inches(6))
        tf = txBox.text_frame
        tf.word_wrap = True

        lines = block.split('\n')
        for i, line in enumerate(lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            run = p.add_run()
            run.text = line
            run.font.size = FONT_SIZE
            run.font.name = FONT_NAME
            run.font.color.rgb = FONT_COLOR
            set_east_asian_font(run, FONT_NAME)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return StreamingResponse(
        buf,
        media_type='application/vnd.openxmlformats-officedocument.presentationml.presentation',
        headers={'Content-Disposition': f'attachment; filename="{pptx_download_filename()}"'},
    )
