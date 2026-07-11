import io
from fastapi import APIRouter
from fastapi.responses import StreamingResponse
from pptx import Presentation
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches
from schemas.line_break import LineBreakEqualsExportRequest
from utils.pptx_helpers import set_east_asian_font
from utils.slide_styles import BIBLE_STYLE, EQUALS_STYLE, SLIDE_HEIGHT, SLIDE_WIDTH
from utils.text_processing import apply_combined_line_break

router = APIRouter(prefix="/line-break/combined/export_ppt", tags=["export"])

ALIGN_MAP = {
    "left": PP_ALIGN.LEFT,
    "center": PP_ALIGN.CENTER,
    "right": PP_ALIGN.RIGHT,
}


@router.post('')
def export_ppt_combined(request: LineBreakEqualsExportRequest):
    entries = apply_combined_line_break(request.text)
    # Each part is its own slide: an outline entry is always one slide, but a
    # multi-verse citation (e.g. '막16:16~17') expands into one slide per verse.
    slides = [(kind, part) for kind, parts in entries for part in parts]
    if not slides:
        slides = [('equals', request.text.strip())]

    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    blank_layout = prs.slide_layouts[6]

    for kind, block in slides:
        style = EQUALS_STYLE if kind == 'equals' else BIBLE_STYLE

        slide = prs.slides.add_slide(blank_layout)
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = style['background_color']

        # Outline slides use a vertically-centered box; Bible slides keep the
        # original top-anchored box, matching each mode's standalone export.
        box_height = SLIDE_HEIGHT - Inches(1) if kind == 'equals' else Inches(6)
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), SLIDE_WIDTH - Inches(1), box_height)
        tf = txBox.text_frame
        tf.word_wrap = True
        if kind == 'equals':
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE

        lines = block.split('\n')
        for i, line in enumerate(lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            if kind == 'equals':
                p.alignment = ALIGN_MAP[request.align]
            run = p.add_run()
            run.text = line
            run.font.size = style['font_size']
            run.font.name = style['font_name']
            if kind == 'equals':
                run.font.bold = True
            run.font.color.rgb = style['font_color']
            set_east_asian_font(run, style['font_name'])

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return StreamingResponse(
        buf,
        media_type='application/vnd.openxmlformats-officedocument.presentationml.presentation',
        headers={'Content-Disposition': 'attachment; filename="output.pptx"'},
    )
