import io
from fastapi import APIRouter
from fastapi.responses import StreamingResponse
from pptx import Presentation
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches
from schemas.line_break import LineBreakEqualsExportRequest
from utils.pptx_helpers import set_east_asian_font
from utils.slide_styles import EQUALS_STYLE, SLIDE_HEIGHT, SLIDE_WIDTH
from utils.text_processing import apply_equals_line_break

router = APIRouter(prefix="/line-break/equals/export_ppt", tags=["export"])

FONT_NAME = EQUALS_STYLE["font_name"]
FONT_SIZE = EQUALS_STYLE["font_size"]
FONT_COLOR = EQUALS_STYLE["font_color"]
BACKGROUND_COLOR = EQUALS_STYLE["background_color"]

ALIGN_MAP = {
    "left": PP_ALIGN.LEFT,
    "center": PP_ALIGN.CENTER,
    "right": PP_ALIGN.RIGHT,
}


@router.post('')
def export_ppt_equals(request: LineBreakEqualsExportRequest):
    align = ALIGN_MAP[request.align]
    text = apply_equals_line_break(request.text)
    blocks = [b.strip() for b in text.split('\n\n') if b.strip()]
    if not blocks:
        blocks = [text.strip()]

    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    blank_layout = prs.slide_layouts[6]

    for block in blocks:
        slide = prs.slides.add_slide(blank_layout)
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = BACKGROUND_COLOR

        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), SLIDE_WIDTH - Inches(1), SLIDE_HEIGHT - Inches(1))
        tf = txBox.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE

        lines = block.split('\n')
        for i, line in enumerate(lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = align
            run = p.add_run()
            run.text = line
            run.font.size = FONT_SIZE
            run.font.name = FONT_NAME
            run.font.bold = True
            run.font.color.rgb = FONT_COLOR
            set_east_asian_font(run, FONT_NAME)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return StreamingResponse(
        buf,
        media_type='application/vnd.openxmlformats-officedocument.presentationml.presentation',
        headers={'Content-Disposition': 'attachment; filename="output.pptx"'},
    )
